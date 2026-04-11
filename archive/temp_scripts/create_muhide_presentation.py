#!/usr/bin/env python3
"""
Create MUHIDE branded AI Sales OS presentation using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors (RGB format)
COLORS = {
    'primary_orange': RGBColor(232, 119, 26),  # #E8771A
    'black': RGBColor(0, 0, 0),
    'dark_black': RGBColor(26, 26, 26),  # #1A1A1A
    'white': RGBColor(255, 255, 255),
    'light_gray': RGBColor(235, 235, 235),  # #EBEBEB
    'dark_gray': RGBColor(74, 74, 74),  # #4A4A4A
    'gray_333': RGBColor(51, 51, 51),  # #333333
}

def add_black_side_bar(slide):
    """Add black vertical bar on left side of content slide"""
    left = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.15), Inches(5.625)
    )
    left.fill.solid()
    left.fill.fore_color.rgb = COLORS['black']
    left.line.color.rgb = COLORS['black']

def add_logo(slide, logo_path, x_offset=8.8, y_offset=0.25):
    """Add MUHIDE logo to slide"""
    try:
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(x_offset), Inches(y_offset), width=Inches(0.9))
    except:
        pass  # Logo not found, skip

def create_title_slide(prs, logo_path):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['black']
    
    # Orange accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.3), Inches(5.625)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['primary_orange']
    bar.line.color.rgb = COLORS['primary_orange']
    
    # Logo
    add_logo(slide, logo_path)
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "MUHIDE"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_orange']
    p.font.name = 'Montserrat'
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(7), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "AI Sales OS Platform"
    p.font.size = Pt(32)
    p.font.color.rgb = COLORS['white']
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(7), Inches(0.4))
    tagline_frame = tagline_box.text_frame
    p = tagline_frame.paragraphs[0]
    p.text = "Trust in Every Transaction•"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['light_gray']
    p.font.italic = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

def create_divider_slide(prs, title):
    """Create divider/section slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['black']
    
    # White gradient effect on right side
    gradient_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5), Inches(0),
        Inches(5), Inches(5.625)
    )
    gradient_box.fill.solid()
    gradient_box.fill.fore_color.rgb = COLORS['white']
    gradient_box.line.color.rgb = COLORS['white']
    
    # Orange bar at top left
    orange_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.2), Inches(0.3)
    )
    orange_bar.fill.solid()
    orange_bar.fill.fore_color.rgb = COLORS['primary_orange']
    orange_bar.line.color.rgb = COLORS['primary_orange']
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title + "•"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = 'Montserrat'
    p.alignment = PP_ALIGN.LEFT

def create_content_slide(prs, title, bullets, logo_path):
    """Create content slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['light_gray']
    
    # Black side bar
    add_black_side_bar(slide)
    
    # Logo
    add_logo(slide, logo_path)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(8.5), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title + "•"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['black']
    p.font.name = 'Montserrat'
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.text = "• " + bullet
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['dark_gray']
        p.font.name = 'Calibri'
        p.space_before = Pt(6)
        p.space_after = Pt(6)

# Main execution
def main():
    # Get paths
    script_dir = os.path.dirname(os.path.abspath(__file__))
    logo_path = os.path.join(script_dir, 'Muhide.png')
    output_path = os.path.join(script_dir, 'AI_Sales_OS_Presentation.pptx')
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Slide 1: Title
    create_title_slide(prs, logo_path)
    
    # Slide 2: Divider
    create_divider_slide(prs, "Executive Summary")
    
    # Slide 3: Market Opportunity
    create_content_slide(prs, "Market Opportunity", [
        "Enterprise sales organizations seek AI-powered solutions",
        "Average sales rep spends 60% time on admin vs selling",
        "$7B+ annual sales tech spending, growing 25% YoY",
        "73% of reps want AI sales tools"
    ], logo_path)
    
    # Slide 4: Divider
    create_divider_slide(prs, "Solution Overview")
    
    # Slide 5: MUHIDE AI Sales OS
    create_content_slide(prs, "MUHIDE AI Sales OS", [
        "Intelligent automation for sales workflows",
        "Real-time opportunity scoring and insights",
        "Seamless CRM integration",
        "ROI-driven metrics and analytics"
    ], logo_path)
    
    # Slide 6: Divider
    create_divider_slide(prs, "Key Features")
    
    # Slide 7: Intelligent Automation
    create_content_slide(prs, "Intelligent Automation", [
        "Lead Qualification: AI-powered lead scoring",
        "Automatic qualification workflows",
        "Sales Engagement: Predictive next actions",
        "Personalized outreach at scale"
    ], logo_path)
    
    # Slide 8: Divider
    create_divider_slide(prs, "Business Case")
    
    # Slide 9: Financial Impact
    create_content_slide(prs, "Financial Impact", [
        "40% Increase in deals closed",
        "3x ROI in first 12 months",
        "25% Reduction in sales cycle",
        "Average rep closes 3-5 additional deals quarterly"
    ], logo_path)
    
    # Slide 10: Divider
    create_divider_slide(prs, "Implementation")
    
    # Slide 11: Deployment Roadmap
    create_content_slide(prs, "Deployment Roadmap", [
        "Phase 1: Assessment (Weeks 1-2) - Sales org analysis",
        "Phase 2: Integration (Weeks 3-6) - CRM setup and AI training",
        "Phase 3: Launch (Weeks 7-8) - Team enablement and go-live"
    ], logo_path)
    
    # Slide 12: Divider
    create_divider_slide(prs, "Competitive Advantage")
    
    # Slide 13: Why MUHIDE
    create_content_slide(prs, "Why MUHIDE", [
        "Purpose-built for enterprise sales",
        "Advanced AI with transparent decision-making",
        "Proven track record with Fortune 500 companies",
        "Unmatched support and customer success"
    ], logo_path)
    
    # Save presentation
    prs.save(output_path)
    print(f"✓ Presentation created: {output_path}")

if __name__ == "__main__":
    main()
