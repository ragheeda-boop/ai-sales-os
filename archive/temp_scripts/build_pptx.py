#!/usr/bin/env python3
"""
Minimal MUHIDE AI Sales OS Presentation Generator
This version uses only standard library + python-pptx to avoid dependency issues
"""

def main():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        import os
    except ImportError as e:
        print(f"Error importing required packages: {e}")
        print("Please install python-pptx: pip install python-pptx")
        return False
    
    # Brand colors (RGB)
    COLORS = {
        'primary_orange': RGBColor(232, 119, 26),
        'black': RGBColor(0, 0, 0),
        'white': RGBColor(255, 255, 255),
        'light_gray': RGBColor(235, 235, 235),
        'dark_gray': RGBColor(74, 74, 74),
    }
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    slide_dir = os.path.dirname(os.path.abspath(__file__))
    logo_path = os.path.join(slide_dir, 'Muhide.png')
    
    # Slide 1: Title
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['black']
    
    # Orange bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(5.625))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['primary_orange']
    bar.line.fill.background()
    
    # Logo
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(8.8), Inches(0.25), width=Inches(0.9))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "MUHIDE"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_orange']
    p.font.name = 'Montserrat'
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(7), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "AI Sales OS Platform"
    p.font.size = Pt(32)
    p.font.color.rgb = COLORS['white']
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(7), Inches(0.4))
    tagline_frame = tagline_box.text_frame
    p = tagline_frame.paragraphs[0]
    p.text = "Trust in Every Transaction•"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['light_gray']
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Divider
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['black']
    
    # White gradient
    gradient = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(0), Inches(5), Inches(5.625))
    gradient.fill.solid()
    gradient.fill.fore_color.rgb = COLORS['white']
    gradient.line.fill.background()
    
    # Orange accent
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(0.3))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['primary_orange']
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "Executive Summary•"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = 'Montserrat'
    p.alignment = PP_ALIGN.LEFT
    
    # Slide 3: Content
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['light_gray']
    
    # Black bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(5.625))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['black']
    bar.line.fill.background()
    
    # Logo
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(8.8), Inches(0.25), width=Inches(0.9))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(8.5), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Market Opportunity•"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['black']
    p.font.name = 'Montserrat'
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    bullets = [
        "Enterprise sales organizations seek AI-powered solutions",
        "Average sales rep spends 60% time on admin vs selling",
        "$7B+ annual sales tech spending, growing 25% YoY",
        "73% of reps want AI sales tools"
    ]
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = "• " + bullet
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['dark_gray']
        p.font.name = 'Calibri'
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Continue with more slides following same pattern...
    # Slide 4-13: Dividers and content slides
    
    slides_data = [
        ("Solution Overview", None),
        ("MUHIDE AI Sales OS", [
            "Intelligent automation for sales workflows",
            "Real-time opportunity scoring and insights",
            "Seamless CRM integration",
            "ROI-driven metrics and analytics"
        ]),
        ("Key Features", None),
        ("Intelligent Automation", [
            "Lead Qualification: AI-powered lead scoring",
            "Automatic qualification workflows",
            "Sales Engagement: Predictive next actions",
            "Personalized outreach at scale"
        ]),
        ("Business Case", None),
        ("Financial Impact", [
            "40% Increase in deals closed",
            "3x ROI in first 12 months",
            "25% Reduction in sales cycle",
            "Average rep closes 3-5 additional deals quarterly"
        ]),
        ("Implementation", None),
        ("Deployment Roadmap", [
            "Phase 1: Assessment (Weeks 1-2)",
            "Phase 2: Integration (Weeks 3-6)",
            "Phase 3: Launch (Weeks 7-8)"
        ]),
        ("Competitive Advantage", None),
        ("Why MUHIDE", [
            "Purpose-built for enterprise sales",
            "Advanced AI with transparent decision-making",
            "Proven track record with Fortune 500 companies",
            "Unmatched support and customer success"
        ]),
    ]
    
    for title, bullets in slides_data:
        if bullets is None:
            # Divider slide
            slide = prs.slides.add_slide(blank_layout)
            background = slide.background
            background.fill.solid()
            background.fill.fore_color.rgb = COLORS['black']
            
            gradient = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(0), Inches(5), Inches(5.625))
            gradient.fill.solid()
            gradient.fill.fore_color.rgb = COLORS['white']
            gradient.line.fill.background()
            
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(0.3))
            accent.fill.solid()
            accent.fill.fore_color.rgb = COLORS['primary_orange']
            accent.line.fill.background()
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            p = title_frame.paragraphs[0]
            p.text = title + "•"
            p.font.size = Pt(54)
            p.font.bold = True
            p.font.color.rgb = COLORS['white']
            p.font.name = 'Montserrat'
        else:
            # Content slide
            slide = prs.slides.add_slide(blank_layout)
            background = slide.background
            background.fill.solid()
            background.fill.fore_color.rgb = COLORS['light_gray']
            
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(5.625))
            bar.fill.solid()
            bar.fill.fore_color.rgb = COLORS['black']
            bar.line.fill.background()
            
            if os.path.exists(logo_path):
                slide.shapes.add_picture(logo_path, Inches(8.8), Inches(0.25), width=Inches(0.9))
            
            title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(8.5), Inches(0.6))
            title_frame = title_box.text_frame
            p = title_frame.paragraphs[0]
            p.text = title + "•"
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = COLORS['black']
            p.font.name = 'Montserrat'
            
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(4))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                p.text = "• " + bullet
                p.level = 0
                p.font.size = Pt(14)
                p.font.color.rgb = COLORS['dark_gray']
                p.font.name = 'Calibri'
                p.space_before = Pt(6)
                p.space_after = Pt(6)
    
    # Save
    output_path = os.path.join(slide_dir, 'AI_Sales_OS_Presentation.pptx')
    prs.save(output_path)
    print(f"✓ Presentation saved: {output_path}")
    return True

if __name__ == "__main__":
    success = main()
    exit(0 if success else 1)
