#!/usr/bin/env python3
"""
AI Sales OS Presentation Generator
Uses python-pptx to create a professional 14-slide presentation
Color palette: MUHIDE brand (orange #E8771A)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Color Palette - MUHIDE Brand
COLORS = {
    'primary': RGBColor(0xE8, 0x77, 0x1A),      # Orange
    'secondary': RGBColor(0x1A, 0x1A, 0x1A),    # Near-black
    'darkBg': RGBColor(0x1A, 0x1A, 0x1A),       # Dark background
    'lightBg': RGBColor(0xF5, 0xF5, 0xF5),      # Light background
    'textLight': RGBColor(0xFF, 0xFF, 0xFF),    # White text
    'textDark': RGBColor(0x1A, 0x1A, 0x1A),     # Dark text
    'accent': RGBColor(0xF0, 0xF0, 0xF0)        # Light gray accent
}

def add_title_slide(prs, title, subtitle, details, notes_text=""):
    """Add a title slide with dark background"""
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['darkBg']

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = COLORS['textLight']
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Details (multiple lines)
    if details:
        details_text = "\n".join(details)
        details_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(2.4))
        details_frame = details_box.text_frame
        details_frame.word_wrap = True
        details_frame.text = details_text
        for para in details_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = COLORS['accent']
            para.alignment = PP_ALIGN.CENTER

    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

    return slide

def add_section_break(prs, title, subtitle):
    """Add a section break slide with dark background"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['darkBg']

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = COLORS['textLight']
    subtitle_para.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title):
    """Add a content slide with light background and orange title bar"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Light background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['lightBg']

    # Title bar with orange background
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.6))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['primary']
    title_shape.line.color.rgb = COLORS['primary']

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['textLight']
    title_para.alignment = PP_ALIGN.LEFT

    return slide

# ============ SLIDE 1: TITLE ============
add_title_slide(
    prs,
    "AI SALES OS",
    "Unified Sales Intelligence Platform",
    [
        "Apollo.io → Python Engine → Notion CRM → GitHub Actions",
        "",
        "44,875 Contacts | 15,407 Companies | 4 Phases | 0 Middleware",
        "",
        "Version 3.1 | March 2026 | Owner: Ragheed"
    ]
)

# ============ SLIDE 2: SECTION BREAK ============
add_section_break(prs, "The Challenge.", "Why we built AI Sales OS")

# ============ SLIDE 3: BEFORE VS AFTER ============
slide = add_content_slide(prs, "Before vs After")

# BEFORE column header
before_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(0.9), Inches(4.2), Inches(0.4))
before_shape.fill.solid()
before_shape.fill.fore_color.rgb = COLORS['secondary']
before_shape.line.color.rgb = COLORS['secondary']

before_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(4.2), Inches(0.3))
before_frame = before_text_box.text_frame
before_frame.text = "BEFORE"
before_para = before_frame.paragraphs[0]
before_para.font.size = Pt(16)
before_para.font.bold = True
before_para.font.color.rgb = COLORS['textLight']
before_para.alignment = PP_ALIGN.CENTER

# BEFORE bullets
before_bullets = [
    "Scattered data across platforms",
    "No automated sync between systems",
    "Manual lead prioritization",
    "Inconsistent scoring criteria",
    "5+ expensive middleware tools"
]

before_text = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(4.0), Inches(3.0))
before_frame = before_text.text_frame
before_frame.word_wrap = True

for i, bullet in enumerate(before_bullets):
    if i == 0:
        p = before_frame.paragraphs[0]
    else:
        p = before_frame.add_paragraph()
    p.text = bullet
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['textDark']
    p.level = 0

# AFTER column header
after_shape = slide.shapes.add_shape(1, Inches(5.3), Inches(0.9), Inches(4.2), Inches(0.4))
after_shape.fill.solid()
after_shape.fill.fore_color.rgb = COLORS['primary']
after_shape.line.color.rgb = COLORS['primary']

after_text_box = slide.shapes.add_textbox(Inches(5.3), Inches(0.95), Inches(4.2), Inches(0.3))
after_frame = after_text_box.text_frame
after_frame.text = "AI SALES OS"
after_para = after_frame.paragraphs[0]
after_para.font.size = Pt(16)
after_para.font.bold = True
after_para.font.color.rgb = COLORS['textLight']
after_para.alignment = PP_ALIGN.CENTER

# AFTER bullets
after_bullets = [
    "Unified data pipeline (Apollo → Notion)",
    "Automatic daily sync (daily_sync.py)",
    "AI lead scoring (0-100 scale)",
    "HOT/WARM/COLD classification",
    "Zero middleware — pure Python + GitHub"
]

after_text = slide.shapes.add_textbox(Inches(5.4), Inches(1.5), Inches(4.0), Inches(3.0))
after_frame = after_text.text_frame
after_frame.word_wrap = True

for i, bullet in enumerate(after_bullets):
    if i == 0:
        p = after_frame.paragraphs[0]
    else:
        p = after_frame.add_paragraph()
    p.text = bullet
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['textDark']
    p.level = 0

# ============ SLIDE 4: SYSTEM ARCHITECTURE ============
slide = add_content_slide(prs, "System Architecture")

# Architecture boxes
boxes = [
    {"x": 0.5, "label": "Apollo.io", "details": "44,875 Contacts\n15,407 Companies"},
    {"x": 2.7, "label": "Python Engine", "details": "daily_sync.py v2.1\nlead_score.py\nauto_tasks.py\n3 Sync Modes"},
    {"x": 4.9, "label": "Notion CRM", "details": "7 Databases\nHOT/WARM/COLD\nDashboards"},
    {"x": 7.1, "label": "GitHub Actions", "details": "Daily Cron\n7:00 AM KSA\n10-Step Pipeline"}
]

for box in boxes:
    # Box shape
    box_shape = slide.shapes.add_shape(1, Inches(box['x']), Inches(1.2), Inches(2.0), Inches(1.8))
    box_shape.fill.solid()
    box_shape.fill.fore_color.rgb = COLORS['secondary']
    box_shape.line.color.rgb = COLORS['primary']
    box_shape.line.width = Pt(2)

    # Label
    label_box = slide.shapes.add_textbox(Inches(box['x']), Inches(1.35), Inches(2.0), Inches(0.35))
    label_frame = label_box.text_frame
    label_frame.text = box['label']
    label_para = label_frame.paragraphs[0]
    label_para.font.size = Pt(14)
    label_para.font.bold = True
    label_para.font.color.rgb = COLORS['primary']
    label_para.alignment = PP_ALIGN.CENTER

    # Details
    details_box = slide.shapes.add_textbox(Inches(box['x'] + 0.1), Inches(1.6), Inches(1.8), Inches(1.2))
    details_frame = details_box.text_frame
    details_frame.text = box['details']
    details_frame.word_wrap = True
    for para in details_frame.paragraphs:
        para.font.size = Pt(10)
        para.font.color.rgb = COLORS['accent']
        para.alignment = PP_ALIGN.CENTER

# Bottom note
note_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.3))
note_frame = note_box.text_frame
note_frame.text = "No middleware — pure Python + GitHub Actions"
note_para = note_frame.paragraphs[0]
note_para.font.size = Pt(12)
note_para.font.bold = True
note_para.font.color.rgb = COLORS['primary']
note_para.alignment = PP_ALIGN.CENTER

# ============ SLIDE 5: DATA AT SCALE ============
slide = add_content_slide(prs, "Data at Scale")

# Big stats
stats = [
    {"num": "44,875", "label": "Contacts"},
    {"num": "15,407", "label": "Companies"},
    {"num": "7", "label": "Databases"},
    {"num": "3", "label": "Sync Modes"}
]

for i, stat in enumerate(stats):
    stat_x = 0.5 + (i * 2.3)

    # Stat box
    stat_shape = slide.shapes.add_shape(1, Inches(stat_x), Inches(1.0), Inches(2.0), Inches(0.9))
    stat_shape.fill.solid()
    stat_shape.fill.fore_color.rgb = COLORS['secondary']
    stat_shape.line.color.rgb = COLORS['primary']
    stat_shape.line.width = Pt(2)

    # Number
    num_box = slide.shapes.add_textbox(Inches(stat_x), Inches(1.15), Inches(2.0), Inches(0.35))
    num_frame = num_box.text_frame
    num_frame.text = stat['num']
    num_para = num_frame.paragraphs[0]
    num_para.font.size = Pt(28)
    num_para.font.bold = True
    num_para.font.color.rgb = COLORS['primary']
    num_para.alignment = PP_ALIGN.CENTER

    # Label
    label_box = slide.shapes.add_textbox(Inches(stat_x), Inches(1.5), Inches(2.0), Inches(0.25))
    label_frame = label_box.text_frame
    label_frame.text = stat['label']
    label_para = label_frame.paragraphs[0]
    label_para.font.size = Pt(11)
    label_para.font.color.rgb = COLORS['accent']
    label_para.alignment = PP_ALIGN.CENTER

# Databases list
db_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(9), Inches(0.3))
db_title_frame = db_title_box.text_frame
db_title_frame.text = "7 Notion Databases"
db_title_para = db_title_frame.paragraphs[0]
db_title_para.font.size = Pt(14)
db_title_para.font.bold = True
db_title_para.font.color.rgb = COLORS['secondary']

databases = [
    "Companies", "Contacts", "Tasks (auto-created by Action Engine)",
    "Opportunities", "Meetings", "Activities", "Email Hub"
]

db_text = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(8.6), Inches(2.7))
db_frame = db_text.text_frame
db_frame.word_wrap = True

for i, db in enumerate(databases):
    if i == 0:
        p = db_frame.paragraphs[0]
    else:
        p = db_frame.add_paragraph()
    p.text = db
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['textDark']
    p.level = 0

# ============ SLIDE 6: SYNC ENGINE ============
slide = add_content_slide(prs, "daily_sync.py v2.1 — Three operational modes")

modes = [
    {"title": "INCREMENTAL", "subtitle": "Daily use", "detail": "Last 26 hours\nOverlap handling\nNo gaps"},
    {"title": "BACKFILL", "subtitle": "Gap recovery", "detail": "Historical range\nCheckpoint resume\nInterruptible"},
    {"title": "FULL", "subtitle": "First-time / rebuild", "detail": "All records\nA-Z partitioning\n2-4 hours"}
]

for i, mode in enumerate(modes):
    mode_x = 0.5 + (i * 3.1)

    # Header box
    header_shape = slide.shapes.add_shape(1, Inches(mode_x), Inches(1.1), Inches(2.9), Inches(0.4))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLORS['primary']
    header_shape.line.color.rgb = COLORS['primary']

    # Title
    title_box = slide.shapes.add_textbox(Inches(mode_x), Inches(1.15), Inches(2.9), Inches(0.3))
    title_frame = title_box.text_frame
    title_frame.text = mode['title']
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(13)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['textLight']
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(mode_x), Inches(1.6), Inches(2.9), Inches(0.25))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = mode['subtitle']
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(10)
    subtitle_para.font.color.rgb = COLORS['secondary']
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Detail
    detail_box = slide.shapes.add_textbox(Inches(mode_x + 0.1), Inches(1.95), Inches(2.7), Inches(1.3))
    detail_frame = detail_box.text_frame
    detail_frame.text = mode['detail']
    detail_frame.word_wrap = True
    for para in detail_frame.paragraphs:
        para.font.size = Pt(9)
        para.font.color.rgb = COLORS['textDark']
        para.alignment = PP_ALIGN.CENTER

# Features
features_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(8.6), Inches(1.5))
features_frame = features_box.text_frame
features_frame.word_wrap = True

feature_list = ["Triple Dedup", "Seniority Normalization", "Safe Boolean Writing", "Pre-loads Notion records"]

for i, feat in enumerate(feature_list):
    if i == 0:
        p = features_frame.paragraphs[0]
    else:
        p = features_frame.add_paragraph()
    p.text = feat
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['textDark']
    p.level = 0

# ============ SLIDE 7: LEAD SCORING ============
slide = add_content_slide(prs, "lead_score.py — Writes Score + Lead Tier")

# Formula title
formula_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.25))
formula_title_frame = formula_title_box.text_frame
formula_title_frame.text = "Scoring Formula (v1.1 — Current Calibrated Weights)"
formula_title_para = formula_title_frame.paragraphs[0]
formula_title_para.font.size = Pt(12)
formula_title_para.font.bold = True
formula_title_para.font.color.rgb = COLORS['secondary']

# Formula box
formula_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(9), Inches(0.55))
formula_shape.fill.solid()
formula_shape.fill.fore_color.rgb = COLORS['secondary']
formula_shape.line.color.rgb = COLORS['primary']
formula_shape.line.width = Pt(1)

formula_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(8.6), Inches(0.35))
formula_frame = formula_box.text_frame
formula_frame.text = "Score = (Intent × 10%) + (Engagement × 10%) + (Company Size × 45%) + (Seniority × 35%)"
formula_frame.word_wrap = True
formula_para = formula_frame.paragraphs[0]
formula_para.font.size = Pt(12)
formula_para.font.color.rgb = COLORS['accent']
formula_para.alignment = PP_ALIGN.CENTER

# Classification title
class_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(9), Inches(0.25))
class_title_frame = class_title_box.text_frame
class_title_frame.text = "Lead Classification"
class_title_para = class_title_frame.paragraphs[0]
class_title_para.font.size = Pt(12)
class_title_para.font.bold = True
class_title_para.font.color.rgb = COLORS['secondary']

# Classification table
table_data = [
    ["HOT", "≥ 80", "Call today — high-value decision makers"],
    ["WARM", "50–79", "Follow up within 48 hours"],
    ["COLD", "< 50", "Monitor only — no action"]
]

rows, cols = 4, 3
left, top, width, height = Inches(0.5), Inches(2.2), Inches(9), Inches(1.6)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Header
for col_idx, header in enumerate(["Tier", "Score", "Action"]):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLORS['primary']

    text_frame = cell.text_frame
    text_frame.paragraphs[0].font.size = Pt(10)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = COLORS['textLight']
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Data rows
for row_idx, row_data in enumerate(table_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = cell_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['lightBg']

        text_frame = cell.text_frame
        text_frame.word_wrap = True
        text_frame.paragraphs[0].font.size = Pt(10)
        text_frame.paragraphs[0].font.color.rgb = COLORS['textDark']
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Note
note_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(0.5))
note_frame = note_box.text_frame
note_frame.text = "Note: v1.1 weights — Size+Seniority weighted high because Intent+Engagement data is sparse. This is expected and correct."
note_frame.word_wrap = True
note_para = note_frame.paragraphs[0]
note_para.font.size = Pt(9)
note_para.font.italic = True
note_para.font.color.rgb = COLORS['secondary']

# ============ SLIDE 8: AUTOMATED PIPELINE ============
slide = add_content_slide(prs, "GitHub Actions — 10-Step Daily Pipeline")

# Steps (two columns)
steps = [
    "1. Checkout repository",
    "2. Setup Python 3.11 with pip cache",
    "3. Install dependencies",
    "4. daily_sync.py (sync last 26 hours)",
    "5. lead_score.py (recalculate + write tier)",
    "6. action_ready_updater.py (5-condition gating)",
    "7. auto_tasks.py (create SLA-based tasks)",
    "8. health_check.py (validate pipeline)",
    "9. Upload logs as artifacts",
    "10. Notify on failure"
]

# Left column
left_steps = "\n".join(steps[:5])
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.5), Inches(2.8))
left_frame = left_box.text_frame
left_frame.text = left_steps
left_frame.word_wrap = True
for para in left_frame.paragraphs:
    para.font.size = Pt(10)
    para.font.color.rgb = COLORS['textDark']

# Right column
right_steps = "\n".join(steps[5:])
right_box = slide.shapes.add_textbox(Inches(5.0), Inches(1.0), Inches(4.5), Inches(2.8))
right_frame = right_box.text_frame
right_frame.text = right_steps
right_frame.word_wrap = True
for para in right_frame.paragraphs:
    para.font.size = Pt(10)
    para.font.color.rgb = COLORS['textDark']

# Cost/timing box
cost_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(4.0), Inches(9), Inches(0.5))
cost_shape.fill.solid()
cost_shape.fill.fore_color.rgb = COLORS['secondary']
cost_shape.line.color.rgb = COLORS['primary']
cost_shape.line.width = Pt(1)

cost_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.1), Inches(8.6), Inches(0.3))
cost_frame = cost_box.text_frame
cost_frame.text = "Schedule: Daily at 7:00 AM KSA (04:00 UTC) | Cost: FREE — 450 min/month of 2,000 free tier"
cost_frame.word_wrap = True
cost_para = cost_frame.paragraphs[0]
cost_para.font.size = Pt(11)
cost_para.font.bold = True
cost_para.font.color.rgb = COLORS['accent']
cost_para.alignment = PP_ALIGN.CENTER

# ============ SLIDE 9: CODE BASE ============
slide = add_content_slide(prs, "Code Base — 10 Active Scripts")

code_data = [
    ["Script", "Purpose", "Status"],
    ["daily_sync.py", "Sync Engine v2.1", "ACTIVE"],
    ["lead_score.py", "Score + Tier Writer", "ACTIVE"],
    ["constants.py", "Unified Field Names", "ACTIVE"],
    ["notion_helpers.py", "Notion API Utils", "ACTIVE"],
    ["auto_tasks.py", "Action Engine (SLA)", "ACTIVE"],
    ["action_ready_updater.py", "5-Condition Gating", "ACTIVE"],
    ["health_check.py", "Pipeline Validator", "ACTIVE"],
    ["webhook_server.py", "Apollo Webhooks", "ACTIVE"],
    ["verify_links.py", "Link Verifier", "ACTIVE"],
    ["job_postings_sync.py", "Job Signals", "PLANNED"]
]

rows, cols = len(code_data), 3
left, top, width, height = Inches(0.3), Inches(0.8), Inches(9.4), Inches(3.5)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Header row
for col_idx, header in enumerate(code_data[0]):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLORS['primary']

    text_frame = cell.text_frame
    text_frame.paragraphs[0].font.size = Pt(9)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = COLORS['textLight']
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Data rows
for row_idx, row_data in enumerate(code_data[1:]):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = cell_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['lightBg']

        text_frame = cell.text_frame
        text_frame.paragraphs[0].font.size = Pt(9)
        text_frame.paragraphs[0].font.color.rgb = COLORS['textDark']
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Tech stack
tech_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.45), Inches(9), Inches(0.3))
tech_frame = tech_box.text_frame
tech_frame.text = "Tech Stack: Python 3.11 | Notion API | Apollo API | GitHub Actions"
tech_para = tech_frame.paragraphs[0]
tech_para.font.size = Pt(11)
tech_para.font.bold = True
tech_para.font.color.rgb = COLORS['secondary']
tech_para.alignment = PP_ALIGN.CENTER

# ============ SLIDE 10: EXECUTION ROADMAP ============
slide = add_content_slide(prs, "Execution Roadmap — 4 Phases")

phases = [
    {"title": "Phase 1\nACTIVATE", "status": "COMPLETE", "color": RGBColor(0x2D, 0x9C, 0x7B)},
    {"title": "Phase 2\nACTION", "status": "CODE COMPLETE", "color": COLORS['primary']},
    {"title": "Phase 3\nENRICH", "status": "NEXT", "color": RGBColor(0x4A, 0x90, 0xE2)},
    {"title": "Phase 4\nOPTIMIZE", "status": "FUTURE", "color": RGBColor(0x99, 0x99, 0x99)}
]

phase_details = [
    ["Full Sync running", "Lead Scoring built", "Seniority normalization", "Safe boolean writing"],
    ["auto_tasks.py built", "action_ready_updater.py", "health_check.py", "Pending first run"],
    ["Job Postings signal", "Job Change detection", "Intent Trend tracking", "Lead Score v2.0"],
    ["Odoo integration", "Revenue tracking", "Analytics", "Full automation"]
]

for idx, phase in enumerate(phases):
    phase_x = 0.4 + idx * 2.35

    # Phase box
    phase_shape = slide.shapes.add_shape(1, Inches(phase_x), Inches(1.2), Inches(2.15), Inches(3.0))
    phase_shape.fill.solid()
    phase_shape.fill.fore_color.rgb = phase['color']
    phase_shape.line.color.rgb = phase['color']

    # Title
    title_box = slide.shapes.add_textbox(Inches(phase_x + 0.1), Inches(1.35), Inches(1.95), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = phase['title']
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(12)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['textLight']
    title_para.alignment = PP_ALIGN.CENTER

    # Status
    status_box = slide.shapes.add_textbox(Inches(phase_x + 0.1), Inches(1.8), Inches(1.95), Inches(0.3))
    status_frame = status_box.text_frame
    status_frame.text = phase['status']
    status_para = status_frame.paragraphs[0]
    status_para.font.size = Pt(9)
    status_para.font.color.rgb = COLORS['accent']
    status_para.alignment = PP_ALIGN.CENTER

    # Details
    details_text = "\n".join(phase_details[idx])
    details_box = slide.shapes.add_textbox(Inches(phase_x + 0.1), Inches(2.2), Inches(1.95), Inches(1.7))
    details_frame = details_box.text_frame
    details_frame.text = details_text
    details_frame.word_wrap = True
    for para in details_frame.paragraphs:
        para.font.size = Pt(8)
        para.font.color.rgb = COLORS['accent']
        para.alignment = PP_ALIGN.CENTER

# Bottom note
note_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.3))
note_frame = note_box.text_frame
note_frame.text = "Overall Progress: Phase 2 code complete — pending first run"
note_para = note_frame.paragraphs[0]
note_para.font.size = Pt(11)
note_para.font.bold = True
note_para.font.color.rgb = COLORS['secondary']
note_para.alignment = PP_ALIGN.CENTER

# ============ SLIDE 11: ACTION ENGINE ============
slide = add_content_slide(prs, "auto_tasks.py + action_ready_updater.py")

# Action Ready conditions title
ar_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.25))
ar_title_frame = ar_title_box.text_frame
ar_title_frame.text = "Action Ready — 5 Gating Conditions"
ar_title_para = ar_title_frame.paragraphs[0]
ar_title_para.font.size = Pt(12)
ar_title_para.font.bold = True
ar_title_para.font.color.rgb = COLORS['secondary']

# Conditions
conditions = [
    "Lead Score ≥ 50",
    "Do Not Call = False",
    "Outreach Status NOT in {Do Not Contact, Bounced, Bad Data}",
    "Stage is NOT 'Customer' or 'Churned'",
    "Has at least one contact method (email or phone)"
]

cond_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(1.3))
cond_frame = cond_text.text_frame
cond_frame.word_wrap = True

for i, cond in enumerate(conditions):
    if i == 0:
        p = cond_frame.paragraphs[0]
    else:
        p = cond_frame.add_paragraph()
    p.text = cond
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS['textDark']
    p.level = 0

# Priority rules title
prio_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.65), Inches(9), Inches(0.25))
prio_title_frame = prio_title_box.text_frame
prio_title_frame.text = "Priority Rules & SLA"
prio_title_para = prio_title_frame.paragraphs[0]
prio_title_para.font.size = Pt(12)
prio_title_para.font.bold = True
prio_title_para.font.color.rgb = COLORS['secondary']

# Priority table
prio_data = [
    ["Tier", "Priority", "Action", "SLA"],
    ["HOT", "Critical", "CALL", "24 hours"],
    ["WARM", "High", "FOLLOW-UP", "48 hours"]
]

rows, cols = len(prio_data), 4
left, top, width, height = Inches(0.5), Inches(3.0), Inches(9), Inches(1.0)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Header
for col_idx, header in enumerate(prio_data[0]):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLORS['primary']

    text_frame = cell.text_frame
    text_frame.paragraphs[0].font.size = Pt(10)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = COLORS['textLight']
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Data
for row_idx, row_data in enumerate(prio_data[1:]):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = cell_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['lightBg']

        text_frame = cell.text_frame
        text_frame.paragraphs[0].font.size = Pt(10)
        text_frame.paragraphs[0].font.color.rgb = COLORS['textDark']
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Features
feat_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.15), Inches(9), Inches(0.25))
feat_title_frame = feat_title_box.text_frame
feat_title_frame.text = "Features:"
feat_title_para = feat_title_frame.paragraphs[0]
feat_title_para.font.size = Pt(12)
feat_title_para.font.bold = True
feat_title_para.font.color.rgb = COLORS['secondary']

feat_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(8.6), Inches(0.8))
feat_frame = feat_box.text_frame
feat_frame.word_wrap = True

features = [
    "Duplicate prevention: checks for existing open tasks before creating",
    "Automated task creation based on Lead Tier and SLA deadlines"
]

for i, feat in enumerate(features):
    if i == 0:
        p = feat_frame.paragraphs[0]
    else:
        p = feat_frame.add_paragraph()
    p.text = feat
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS['textDark']
    p.level = 0

# ============ SLIDE 12: SIGNALS & ENRICHMENT ============
slide = add_content_slide(prs, "Phase 3 — Signals & Enrichment")

# Coming in Phase 3
coming_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.25))
coming_title_frame = coming_title_box.text_frame
coming_title_frame.text = "Coming in Phase 3"
coming_title_para = coming_title_frame.paragraphs[0]
coming_title_para.font.size = Pt(12)
coming_title_para.font.bold = True
coming_title_para.font.color.rgb = COLORS['secondary']

coming_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(1.2))
coming_frame = coming_box.text_frame
coming_frame.word_wrap = True

coming_items = [
    "Job Postings Signal — detect hiring activity at target accounts",
    "Job Change Detection — identify when decision-makers change roles",
    "Intent Trend Tracking — compare intent scores between syncs"
]

for i, item in enumerate(coming_items):
    if i == 0:
        p = coming_frame.paragraphs[0]
    else:
        p = coming_frame.add_paragraph()
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['textDark']
    p.level = 0

# Lead Score v2.0 title
v2_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(9), Inches(0.25))
v2_title_frame = v2_title_box.text_frame
v2_title_frame.text = "Lead Score v2.0 Weights (activate with signals data)"
v2_title_para = v2_title_frame.paragraphs[0]
v2_title_para.font.size = Pt(12)
v2_title_para.font.bold = True
v2_title_para.font.color.rgb = COLORS['secondary']

# v2.0 formula box
v2_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(3.0), Inches(9), Inches(0.5))
v2_shape.fill.solid()
v2_shape.fill.fore_color.rgb = COLORS['secondary']
v2_shape.line.color.rgb = COLORS['primary']
v2_shape.line.width = Pt(1)

v2_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.1), Inches(8.6), Inches(0.3))
v2_frame = v2_box.text_frame
v2_frame.text = "Intent 30% + Engagement 25% + Signals 25% + Company Size 10% + Seniority 10%"
v2_frame.word_wrap = True
v2_para = v2_frame.paragraphs[0]
v2_para.font.size = Pt(11)
v2_para.font.color.rgb = COLORS['accent']
v2_para.alignment = PP_ALIGN.CENTER

# Note
note_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.6))
note_frame = note_box.text_frame
note_frame.text = "Note: DO NOT activate v2.0 until signals data is actually populated. Empty fields would weaken all scores."
note_frame.word_wrap = True
note_para = note_frame.paragraphs[0]
note_para.font.size = Pt(10)
note_para.font.italic = True
note_para.font.color.rgb = COLORS['secondary']

# ============ SLIDE 13: CURRENT STATUS ============
slide = add_content_slide(prs, "Current Status")

status_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(4.0))
status_frame = status_box.text_frame
status_frame.word_wrap = True

status_items = [
    ("Phase 1: ACTIVATE", True),
    ("✓ COMPLETE — Full sync, lead scoring, seniority normalization, safe booleans", False),
    ("", False),
    ("Phase 2: ACTION", True),
    ("✓ CODE COMPLETE — auto_tasks.py, action_ready_updater.py, health_check.py", False),
    ("", False),
    ("Remaining Tasks:", True),
    ("1. Run lead_score.py --force (write Lead Tier for all contacts)", False),
    ("2. Run action_ready_updater.py (evaluate 5 conditions)", False),
    ("3. Run auto_tasks.py --dry-run (validate task creation)", False),
    ("4. Push code to GitHub, add NOTION_DATABASE_ID_TASKS secret", False),
    ("5. Activate daily pipeline (GitHub Actions)", False)
]

for i, (text, is_bold) in enumerate(status_items):
    if i == 0:
        p = status_frame.paragraphs[0]
    else:
        p = status_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(10)
    p.font.bold = is_bold
    p.font.color.rgb = COLORS['textDark']
    p.level = 0

# ============ SLIDE 14: WHAT'S NEXT ============
add_title_slide(
    prs,
    "What's Next",
    "",
    [
        "1. Run lead_score.py --force (write Lead Tier for all contacts)",
        "2. Run action_ready_updater.py + auto_tasks.py --dry-run",
        "3. Push code to GitHub, add NOTION_DATABASE_ID_TASKS secret",
        "4. Activate daily pipeline (GitHub Actions)",
        "5. Phase 3: Job Postings + Intent Signals",
        "",
        "AI SALES OS — Built by Ragheed | Version 3.1 | March 2026"
    ]
)

# Save presentation
output_path = "/sessions/wizardly-elegant-einstein/mnt/AI Sales OS/AI_Sales_OS_Presentation.pptx"
prs.save(output_path)
print(f"✓ Presentation created successfully!")
print(f"✓ File saved to: {output_path}")
